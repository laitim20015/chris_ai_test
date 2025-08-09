"""
測試數據生成器

生成各種格式的測試文檔和數據，用於測試系統功能。
"""

import json
import random
from pathlib import Path
from typing import List, Dict, Any
from datetime import datetime, timedelta

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.colors import black, blue, red, green
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
import io
from PIL import Image as PILImage, ImageDraw


class TestDataGenerator:
    """測試數據生成器"""
    
    def __init__(self, output_dir: Path):
        """
        初始化測試數據生成器
        
        Args:
            output_dir: 輸出目錄
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # 創建子目錄
        (self.output_dir / "documents").mkdir(exist_ok=True)
        (self.output_dir / "images").mkdir(exist_ok=True)
        (self.output_dir / "expected_outputs").mkdir(exist_ok=True)
    
    def generate_test_images(self) -> List[Path]:
        """生成測試圖片"""
        image_paths = []
        
        # 生成圖表樣本
        chart_path = self.output_dir / "images" / "chart_sample.jpg"
        self._create_chart_image(chart_path)
        image_paths.append(chart_path)
        
        # 生成圖示樣本
        diagram_path = self.output_dir / "images" / "diagram_sample.png"
        self._create_diagram_image(diagram_path)
        image_paths.append(diagram_path)
        
        # 生成表格圖片
        table_path = self.output_dir / "images" / "table_sample.png"
        self._create_table_image(table_path)
        image_paths.append(table_path)
        
        return image_paths
    
    def _create_chart_image(self, path: Path):
        """創建圖表圖片"""
        img = PILImage.new('RGB', (400, 300), 'white')
        draw = ImageDraw.Draw(img)
        
        # 繪製簡單的柱狀圖
        bars = [(50, 200), (100, 150), (150, 180), (200, 120)]
        colors = ['blue', 'red', 'green', 'orange']
        
        for i, ((x, height), color) in enumerate(zip(bars, colors)):
            x_pos = 50 + i * 80
            draw.rectangle([x_pos, 250 - height, x_pos + 60, 250], fill=color)
            draw.text((x_pos + 20, 260), f"Q{i+1}", fill='black')
        
        # 添加標題
        draw.text((150, 20), "Sales Chart", fill='black')
        draw.text((20, 50), "Revenue", fill='black')
        
        img.save(path, 'JPEG')
    
    def _create_diagram_image(self, path: Path):
        """創建圖示圖片"""
        img = PILImage.new('RGB', (350, 250), 'white')
        draw = ImageDraw.Draw(img)
        
        # 繪製流程圖
        boxes = [(50, 50, 150, 90), (200, 50, 300, 90), (125, 150, 225, 190)]
        labels = ["Process A", "Process B", "Result"]
        
        for (x1, y1, x2, y2), label in zip(boxes, labels):
            draw.rectangle([x1, y1, x2, y2], outline='black', fill='lightblue')
            text_x = x1 + (x2 - x1) // 2 - len(label) * 3
            text_y = y1 + (y2 - y1) // 2 - 5
            draw.text((text_x, text_y), label, fill='black')
        
        # 繪製箭頭
        draw.line((150, 70, 200, 70), fill='black', width=2)
        draw.line((175, 90, 175, 150), fill='black', width=2)
        
        img.save(path, 'PNG')
    
    def _create_table_image(self, path: Path):
        """創建表格圖片"""
        img = PILImage.new('RGB', (400, 200), 'white')
        draw = ImageDraw.Draw(img)
        
        # 繪製表格
        rows, cols = 4, 3
        cell_width, cell_height = 120, 40
        start_x, start_y = 20, 20
        
        # 表格標題
        headers = ["Product", "Price", "Stock"]
        data = [
            ["Widget A", "$10.99", "50"],
            ["Widget B", "$15.99", "30"],
            ["Widget C", "$8.99", "75"]
        ]
        
        # 繪製表格框架
        for i in range(rows + 1):
            y = start_y + i * cell_height
            draw.line((start_x, y, start_x + cols * cell_width, y), fill='black')
        
        for j in range(cols + 1):
            x = start_x + j * cell_width
            draw.line((x, start_y, x, start_y + rows * cell_height), fill='black')
        
        # 填充內容
        for j, header in enumerate(headers):
            x = start_x + j * cell_width + 10
            draw.text((x, start_y + 10), header, fill='black')
        
        for i, row in enumerate(data):
            for j, cell in enumerate(row):
                x = start_x + j * cell_width + 10
                y = start_y + (i + 1) * cell_height + 10
                draw.text((x, y), cell, fill='black')
        
        img.save(path, 'PNG')
    
    def generate_test_pdf(self, filename: str = "sample.pdf") -> Path:
        """生成測試PDF文檔"""
        pdf_path = self.output_dir / "documents" / filename
        
        # 創建PDF文檔
        doc = SimpleDocTemplate(str(pdf_path), pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # 標題
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            textColor=blue,
            alignment=1  # 居中
        )
        story.append(Paragraph("測試文檔標題", title_style))
        story.append(Spacer(1, 12))
        
        # 正文段落
        normal_style = styles['Normal']
        paragraphs = [
            "這是第一段測試內容。本段落用於測試文本提取功能，包含了基本的中文文字處理。",
            "圖1顯示了系統的整體架構。通過這個圖表，我們可以清楚地看到各個組件之間的關係。",
            "第三段內容涉及數據分析。表1中的數據來源於實際的業務場景，具有很高的參考價值。",
            "如圖2所示，性能測試結果表明系統在高負載下仍能保持穩定運行。",
            "最後一段是總結內容。通過以上分析，我們可以得出結論：系統設計合理，性能優良。"
        ]
        
        for i, para_text in enumerate(paragraphs):
            story.append(Paragraph(para_text, normal_style))
            story.append(Spacer(1, 12))
            
            # 在某些段落後添加圖片
            if i == 1:  # 在"圖1"之後
                img_path = self.output_dir / "images" / "chart_sample.jpg"
                if img_path.exists():
                    story.append(Image(str(img_path), width=3*inch, height=2.25*inch))
                    story.append(Spacer(1, 12))
            
            elif i == 3:  # 在"圖2"之後
                img_path = self.output_dir / "images" / "diagram_sample.png"
                if img_path.exists():
                    story.append(Image(str(img_path), width=2.5*inch, height=1.875*inch))
                    story.append(Spacer(1, 12))
        
        # 添加表格
        table_data = [
            ['項目', '數值', '單位'],
            ['性能指標A', '95.5', '%'],
            ['性能指標B', '1.2', '秒'],
            ['性能指標C', '256', 'MB']
        ]
        
        table = Table(table_data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), blue),
            ('TEXTCOLOR', (0, 0), (-1, 0), 'white'),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), 'beige'),
            ('GRID', (0, 0), (-1, -1), 1, black)
        ]))
        
        story.append(Paragraph("表1：性能指標統計", styles['Heading2']))
        story.append(Spacer(1, 6))
        story.append(table)
        
        # 構建PDF
        doc.build(story)
        return pdf_path
    
    def generate_test_docx(self, filename: str = "sample.docx") -> Path:
        """生成測試Word文檔"""
        try:
            from docx import Document
            from docx.shared import Inches
        except ImportError:
            # 如果沒有python-docx，創建一個模擬文件
            docx_path = self.output_dir / "documents" / filename
            with open(docx_path, "wb") as f:
                f.write(b"PK\x03\x04")  # ZIP頭部
                f.write(b"mock docx content for testing")
            return docx_path
        
        docx_path = self.output_dir / "documents" / filename
        doc = Document()
        
        # 添加標題
        doc.add_heading('測試Word文檔', 0)
        
        # 添加段落
        paragraphs = [
            "這是Word文檔的第一段內容。用於測試Word解析器的文本提取能力。",
            "第二段提到了圖片。下方的圖1展示了系統架構圖。",
            "第三段包含表格引用。詳細數據請參考表1。",
            "第四段是關於性能分析的內容。圖2顯示了性能測試結果。",
            "最後一段是總結性內容。"
        ]
        
        for para_text in paragraphs:
            doc.add_paragraph(para_text)
        
        # 添加圖片（如果存在）
        img_path = self.output_dir / "images" / "chart_sample.jpg"
        if img_path.exists():
            doc.add_paragraph("圖1：系統架構圖")
            doc.add_picture(str(img_path), width=Inches(4))
        
        # 添加表格
        table = doc.add_table(rows=1, cols=3)
        table.style = 'Light Grid Accent 1'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = '項目'
        hdr_cells[1].text = '數值'
        hdr_cells[2].text = '說明'
        
        data = [
            ('指標A', '95.5%', '性能良好'),
            ('指標B', '1.2秒', '響應快速'),
            ('指標C', '256MB', '內存適中')
        ]
        
        for item, value, desc in data:
            row_cells = table.add_row().cells
            row_cells[0].text = item
            row_cells[1].text = value
            row_cells[2].text = desc
        
        doc.save(str(docx_path))
        return docx_path
    
    def generate_test_pptx(self, filename: str = "sample.pptx") -> Path:
        """生成測試PowerPoint文檔"""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            # 如果沒有python-pptx，創建一個模擬文件
            pptx_path = self.output_dir / "documents" / filename
            with open(pptx_path, "wb") as f:
                f.write(b"PK\x03\x04")  # ZIP頭部
                f.write(b"mock pptx content for testing")
            return pptx_path
        
        pptx_path = self.output_dir / "documents" / filename
        prs = Presentation()
        
        # 第一張幻燈片：標題頁
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "測試演示文稿"
        subtitle.text = "用於測試PowerPoint解析器功能"
        
        # 第二張幻燈片：內容頁
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = '主要內容'
        
        tf = body_shape.text_frame
        tf.text = '系統功能概述'
        
        p = tf.add_paragraph()
        p.text = '圖1顯示了系統架構'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = '表1包含了性能數據'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = '圖2展示了測試結果'
        p.level = 1
        
        # 第三張幻燈片：圖片頁
        if (self.output_dir / "images" / "chart_sample.jpg").exists():
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 添加標題
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = "圖1：性能圖表"
            
            # 添加圖片
            img_path = str(self.output_dir / "images" / "chart_sample.jpg")
            slide.shapes.add_picture(img_path, Inches(2), Inches(2), Inches(4), Inches(3))
        
        prs.save(str(pptx_path))
        return pptx_path
    
    def generate_expected_outputs(self) -> Dict[str, Path]:
        """生成預期輸出結果"""
        outputs = {}
        
        # 生成預期的Markdown輸出
        md_path = self.output_dir / "expected_outputs" / "sample_output.md"
        with open(md_path, "w", encoding="utf-8") as f:
            f.write("""# 測試文檔標題

這是第一段測試內容。本段落用於測試文本提取功能，包含了基本的中文文字處理。

圖1顯示了系統的整體架構。通過這個圖表，我們可以清楚地看到各個組件之間的關係。

![圖1](images/chart_sample.jpg)

第三段內容涉及數據分析。表1中的數據來源於實際的業務場景，具有很高的參考價值。

如圖2所示，性能測試結果表明系統在高負載下仍能保持穩定運行。

![圖2](images/diagram_sample.png)

最後一段是總結內容。通過以上分析，我們可以得出結論：系統設計合理，性能優良。

## 表1：性能指標統計

| 項目 | 數值 | 單位 |
|------|------|------|
| 性能指標A | 95.5 | % |
| 性能指標B | 1.2 | 秒 |
| 性能指標C | 256 | MB |
""")
        outputs["markdown"] = md_path
        
        # 生成預期的關聯結果
        associations_path = self.output_dir / "expected_outputs" / "associations.json"
        associations_data = {
            "document_id": "test_doc_001",
            "associations": [
                {
                    "text_block_id": "text_002",
                    "image_id": "img_001",
                    "caption_score": 0.95,
                    "spatial_score": 0.85,
                    "semantic_score": 0.75,
                    "final_score": 0.88,
                    "association_type": "direct_reference",
                    "confidence": 0.92
                },
                {
                    "text_block_id": "text_004",
                    "image_id": "img_002",
                    "caption_score": 0.90,
                    "spatial_score": 0.80,
                    "semantic_score": 0.70,
                    "final_score": 0.83,
                    "association_type": "direct_reference",
                    "confidence": 0.87
                }
            ],
            "statistics": {
                "total_text_blocks": 5,
                "total_images": 2,
                "total_associations": 2,
                "high_confidence_associations": 2,
                "average_confidence": 0.895
            }
        }
        
        with open(associations_path, "w", encoding="utf-8") as f:
            json.dump(associations_data, f, indent=2, ensure_ascii=False)
        outputs["associations"] = associations_path
        
        return outputs
    
    def generate_performance_test_data(self, num_docs: int = 10) -> List[Path]:
        """生成性能測試數據"""
        docs = []
        
        for i in range(num_docs):
            # 生成不同大小的文檔
            if i < 3:  # 小文檔
                filename = f"small_doc_{i+1}.pdf"
                doc_path = self.generate_simple_pdf(filename, pages=1, paragraphs_per_page=5)
            elif i < 7:  # 中等文檔
                filename = f"medium_doc_{i+1}.pdf"
                doc_path = self.generate_simple_pdf(filename, pages=5, paragraphs_per_page=10)
            else:  # 大文檔
                filename = f"large_doc_{i+1}.pdf"
                doc_path = self.generate_simple_pdf(filename, pages=20, paragraphs_per_page=15)
            
            docs.append(doc_path)
        
        return docs
    
    def generate_simple_pdf(self, filename: str, pages: int = 1, paragraphs_per_page: int = 5) -> Path:
        """生成簡單的PDF文檔"""
        pdf_path = self.output_dir / "documents" / filename
        
        c = canvas.Canvas(str(pdf_path), pagesize=letter)
        width, height = letter
        
        for page_num in range(pages):
            y_position = height - 50
            
            # 頁面標題
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, y_position, f"第{page_num + 1}頁測試內容")
            y_position -= 40
            
            # 段落內容
            c.setFont("Helvetica", 12)
            for para_num in range(paragraphs_per_page):
                if y_position < 100:  # 接近頁面底部
                    break
                
                para_text = f"這是第{para_num + 1}段測試文字。" * random.randint(2, 5)
                
                # 簡單的文字換行
                words = para_text.split()
                line = ""
                for word in words:
                    if len(line + word) < 80:  # 簡單的行長度控制
                        line += word + " "
                    else:
                        c.drawString(50, y_position, line)
                        y_position -= 20
                        line = word + " "
                        
                        if y_position < 100:
                            break
                
                if line and y_position >= 100:
                    c.drawString(50, y_position, line)
                    y_position -= 30
            
            if page_num < pages - 1:
                c.showPage()
        
        c.save()
        return pdf_path
    
    def cleanup_test_data(self):
        """清理測試數據"""
        import shutil
        if self.output_dir.exists():
            shutil.rmtree(self.output_dir)


def generate_all_test_data(output_dir: str = "tests/fixtures") -> Dict[str, Any]:
    """生成所有測試數據的便捷函數"""
    generator = TestDataGenerator(Path(output_dir))
    
    print("🖼️ 生成測試圖片...")
    images = generator.generate_test_images()
    
    print("📄 生成測試PDF...")
    pdf_path = generator.generate_test_pdf()
    
    print("📝 生成測試Word文檔...")
    docx_path = generator.generate_test_docx()
    
    print("🎯 生成測試PowerPoint...")
    pptx_path = generator.generate_test_pptx()
    
    print("📊 生成預期輸出...")
    expected_outputs = generator.generate_expected_outputs()
    
    print("⚡ 生成性能測試數據...")
    performance_docs = generator.generate_performance_test_data()
    
    print("✅ 測試數據生成完成!")
    
    return {
        "images": images,
        "documents": {
            "pdf": pdf_path,
            "docx": docx_path,
            "pptx": pptx_path
        },
        "expected_outputs": expected_outputs,
        "performance_docs": performance_docs,
        "output_dir": generator.output_dir
    }


if __name__ == "__main__":
    # 直接運行時生成測試數據
    result = generate_all_test_data()
    print(f"📁 測試數據已生成到: {result['output_dir']}")
    
    print("\n📋 生成的文件:")
    for category, files in result.items():
        if category != "output_dir":
            print(f"  {category}: {files}")
