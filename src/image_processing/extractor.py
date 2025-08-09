"""
圖片提取器

負責從各種文檔格式中提取圖片，並進行初步處理
遵循項目規格文件中的圖片處理流程：
原始圖片 → 格式標準化 → 尺寸優化 → 質量壓縮 → 命名規範 → 存儲上傳
"""

import os
import io
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
import fitz  # PyMuPDF
from PIL import Image
import base64

from src.config.logging_config import get_logger
from src.parsers.base import ImageContent, BoundingBox

logger = get_logger(__name__)

@dataclass
class ExtractedImage:
    """提取的圖片數據結構"""
    image_data: bytes                    # 圖片二進制數據
    format: str                         # 圖片格式 (jpg, png, etc.)
    size: Tuple[int, int]              # 尺寸 (width, height)
    source_page: int                   # 來源頁面
    source_position: BoundingBox       # 原始位置
    hash: str                          # 圖片內容哈希
    filename: str                      # 建議文件名
    metadata: Dict[str, Any]           # 額外元數據

class ImageExtractor:
    """
    圖片提取器
    
    按照規格文件要求，實現：
    1. 多格式圖片提取
    2. 格式標準化處理
    3. 去重檢測
    4. 命名規範實施
    """
    
    def __init__(self, 
                 supported_formats: Optional[List[str]] = None,
                 min_size: Tuple[int, int] = (50, 50),
                 max_size: Tuple[int, int] = (4096, 4096)):
        """
        初始化圖片提取器
        
        Args:
            supported_formats: 支持的圖片格式
            min_size: 最小圖片尺寸 (width, height)
            max_size: 最大圖片尺寸 (width, height)
        """
        self.supported_formats = supported_formats or ['PNG', 'JPEG', 'JPG', 'GIF', 'BMP', 'TIFF']
        self.min_size = min_size
        self.max_size = max_size
        self.extracted_hashes = set()  # 去重用的哈希集合
        
        logger.info(f"圖片提取器初始化完成 - 支持格式: {self.supported_formats}")
    
    def extract_from_pdf(self, pdf_path: str, document_name: str) -> List[ExtractedImage]:
        """
        從PDF文檔提取圖片
        
        Args:
            pdf_path: PDF文件路徑
            document_name: 文檔名稱（用於命名）
            
        Returns:
            提取的圖片列表
        """
        logger.info(f"開始從PDF提取圖片: {pdf_path}")
        extracted_images = []
        
        try:
            doc = fitz.open(pdf_path)
            logger.info(f"PDF文檔已打開，共{len(doc)}頁")
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # 提取嵌入圖片
                images = page.get_images(full=True)
                logger.debug(f"頁面{page_num + 1}發現{len(images)}個嵌入圖片")
                
                for img_index, img in enumerate(images):
                    try:
                        logger.debug(f"處理圖片 - 頁面{page_num + 1}, 圖片{img_index}, xref={img[0]}")
                        extracted_img = self._extract_pdf_image(
                            doc, page, img, page_num + 1, img_index, document_name
                        )
                        if extracted_img:
                            extracted_images.append(extracted_img)
                            logger.debug(f"成功提取圖片: {extracted_img.filename}")
                        else:
                            logger.debug(f"圖片提取返回None")
                    except Exception as e:
                        logger.warning(f"提取PDF圖片失敗 - 頁面{page_num + 1}, 圖片{img_index}: {e}")
                        import traceback
                        logger.debug(f"錯誤詳情: {traceback.format_exc()}")
                        continue
                
                # 提取矢量圖形（轉為光柵圖片）
                vector_images = self._extract_vector_graphics(page, page_num + 1, document_name)
                logger.debug(f"頁面{page_num + 1}發現{len(vector_images)}個矢量圖形")
                extracted_images.extend(vector_images)
            
            doc.close()
            
        except Exception as e:
            logger.error(f"PDF圖片提取失敗: {e}")
            raise
        
        logger.info(f"PDF圖片提取完成 - 提取了{len(extracted_images)}張圖片")
        return extracted_images
    
    def _extract_pdf_image(self, doc: fitz.Document, page: fitz.Page, 
                          img: tuple, page_num: int, img_index: int, 
                          document_name: str) -> Optional[ExtractedImage]:
        """提取PDF中的單個嵌入圖片"""
        
        xref = img[0]
        pix = fitz.Pixmap(doc, xref)
        
        # 檢查圖片格式和尺寸
        if pix.n - pix.alpha < 3:  # 灰度或單色圖片
            pix = fitz.Pixmap(fitz.csRGB, pix)
        
        logger.debug(f"圖片尺寸檢查: {pix.width}x{pix.height}, 最小要求: {self.min_size}")
        
        if pix.width < self.min_size[0] or pix.height < self.min_size[1]:
            logger.debug(f"圖片尺寸太小，跳過: {pix.width}x{pix.height}")
            pix = None
            return None
        
        # 轉換為PIL Image進行進一步處理
        img_data = pix.tobytes("png")
        pil_image = Image.open(io.BytesIO(img_data))
        
        # 生成圖片哈希用於去重
        image_hash = self._calculate_image_hash(img_data)
        if image_hash in self.extracted_hashes:
            logger.debug(f"跳過重複圖片 - 頁面{page_num}, 索引{img_index}")
            pix = None
            return None
        
        self.extracted_hashes.add(image_hash)
        
        # 獲取圖片在頁面中的位置
        img_rects = page.get_image_rects(xref)
        bbox = BoundingBox(0, 0, pix.width, pix.height)
        if img_rects:
            rect = img_rects[0]
            bbox = BoundingBox(rect.x0, rect.y0, rect.x1, rect.y1)
        
        # 按照規格文件的命名規範生成文件名
        # {文件名}_{頁碼}_{圖片序號}_{時間戳}.{格式}
        import time
        timestamp = int(time.time())
        filename = f"{document_name}_p{page_num:03d}_img{img_index:03d}_{timestamp}.png"
        
        # 創建提取結果
        extracted_image = ExtractedImage(
            image_data=img_data,
            format="PNG",
            size=(pix.width, pix.height),
            source_page=page_num,
            source_position=bbox,
            hash=image_hash,
            filename=filename,
            metadata={
                "source_type": "embedded",
                "original_format": "PDF_embedded",
                "samples_per_pixel": pix.n,
                "color_space": pix.colorspace.name if pix.colorspace else "unknown"
            }
        )
        
        pix = None
        return extracted_image
    
    def _extract_vector_graphics(self, page: fitz.Page, page_num: int, 
                                document_name: str) -> List[ExtractedImage]:
        """
        提取矢量圖形並轉換為光柵圖片
        
        Args:
            page: PDF頁面對象
            page_num: 頁面號碼
            document_name: 文檔名稱
            
        Returns:
            轉換後的圖片列表
        """
        vector_images = []
        
        try:
            # 獲取頁面中的繪圖對象
            drawings = page.get_drawings()
            
            if not drawings:
                return vector_images
            
            # 將矢量圖形分組
            grouped_drawings = self._group_drawings_by_proximity(drawings)
            
            for group_index, drawing_group in enumerate(grouped_drawings):
                try:
                    # 計算包圍盒
                    bbox = self._calculate_drawings_bbox(drawing_group)
                    
                    if bbox.width < self.min_size[0] or bbox.height < self.min_size[1]:
                        continue
                    
                    # 渲染矢量圖形為圖片
                    vector_image = self._render_vector_to_image(
                        page, bbox, page_num, group_index, document_name
                    )
                    
                    if vector_image:
                        vector_images.append(vector_image)
                        
                except Exception as e:
                    logger.warning(f"矢量圖形處理失敗 - 頁面{page_num}, 組{group_index}: {e}")
                    continue
        
        except Exception as e:
            logger.warning(f"矢量圖形提取失敗 - 頁面{page_num}: {e}")
        
        return vector_images
    
    def _group_drawings_by_proximity(self, drawings: List[Dict], 
                                   proximity_threshold: float = 50.0) -> List[List[Dict]]:
        """
        按照空間接近程度對繪圖對象進行分組
        
        Args:
            drawings: 繪圖對象列表
            proximity_threshold: 接近程度閾值
            
        Returns:
            分組後的繪圖對象列表
        """
        if not drawings:
            return []
        
        groups = []
        processed = set()
        
        for i, drawing in enumerate(drawings):
            if i in processed:
                continue
            
            # 創建新組
            current_group = [drawing]
            processed.add(i)
            
            # 查找接近的繪圖對象
            for j, other_drawing in enumerate(drawings):
                if j in processed:
                    continue
                
                if self._are_drawings_close(drawing, other_drawing, proximity_threshold):
                    current_group.append(other_drawing)
                    processed.add(j)
            
            groups.append(current_group)
        
        return groups
    
    def _are_drawings_close(self, drawing1: Dict, drawing2: Dict, 
                           threshold: float) -> bool:
        """檢查兩個繪圖對象是否接近"""
        try:
            rect1 = drawing1.get('rect', fitz.Rect(0, 0, 0, 0))
            rect2 = drawing2.get('rect', fitz.Rect(0, 0, 0, 0))
            
            # 計算中心點距離
            center1 = ((rect1.x0 + rect1.x1) / 2, (rect1.y0 + rect1.y1) / 2)
            center2 = ((rect2.x0 + rect2.x1) / 2, (rect2.y0 + rect2.y1) / 2)
            
            distance = ((center1[0] - center2[0])**2 + (center1[1] - center2[1])**2)**0.5
            
            return distance <= threshold
        except:
            return False
    
    def _calculate_drawings_bbox(self, drawings: List[Dict]) -> BoundingBox:
        """計算繪圖對象組的包圍盒"""
        if not drawings:
            return BoundingBox(0, 0, 0, 0)
        
        min_x = min_y = float('inf')
        max_x = max_y = float('-inf')
        
        for drawing in drawings:
            rect = drawing.get('rect', fitz.Rect(0, 0, 0, 0))
            min_x = min(min_x, rect.x0)
            min_y = min(min_y, rect.y0)
            max_x = max(max_x, rect.x1)
            max_y = max(max_y, rect.y1)
        
        return BoundingBox(min_x, min_y, max_x, max_y)
    
    def _render_vector_to_image(self, page: fitz.Page, bbox: BoundingBox,
                               page_num: int, group_index: int, 
                               document_name: str) -> Optional[ExtractedImage]:
        """將矢量圖形渲染為光柵圖片"""
        
        try:
            # 創建裁切區域
            clip_rect = fitz.Rect(bbox.x1, bbox.y1, bbox.x2, bbox.y2)
            
            # 渲染頁面片段
            matrix = fitz.Matrix(2.0, 2.0)  # 2倍分辨率
            pix = page.get_pixmap(matrix=matrix, clip=clip_rect)
            
            if pix.width < self.min_size[0] or pix.height < self.min_size[1]:
                pix = None
                return None
            
            # 轉換為PNG格式
            img_data = pix.tobytes("png")
            
            # 生成哈希並檢查重複
            image_hash = self._calculate_image_hash(img_data)
            if image_hash in self.extracted_hashes:
                pix = None
                return None
            
            self.extracted_hashes.add(image_hash)
            
            # 生成文件名
            import time
            timestamp = int(time.time())
            filename = f"{document_name}_p{page_num:03d}_vector{group_index:03d}_{timestamp}.png"
            
            # 創建提取結果
            extracted_image = ExtractedImage(
                image_data=img_data,
                format="PNG",
                size=(pix.width, pix.height),
                source_page=page_num,
                source_position=bbox,
                hash=image_hash,
                filename=filename,
                metadata={
                    "source_type": "vector_graphics",
                    "original_format": "PDF_vector",
                    "render_scale": 2.0,
                    "group_index": group_index
                }
            )
            
            pix = None
            return extracted_image
            
        except Exception as e:
            logger.warning(f"矢量圖形渲染失敗: {e}")
            return None
    
    def extract_from_word(self, word_path: str, document_name: str) -> List[ExtractedImage]:
        """
        從Word文檔提取圖片
        
        Args:
            word_path: Word文件路徑
            document_name: 文檔名稱
            
        Returns:
            提取的圖片列表
        """
        logger.info(f"開始從Word文檔提取圖片: {word_path}")
        extracted_images = []
        
        try:
            from docx import Document
            from docx.image.exceptions import UnrecognizedImageError
            
            doc = Document(word_path)
            
            # 遍歷文檔中的所有關係來查找圖片
            img_index = 0
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        # 獲取圖片數據
                        image_data = rel.target_part.blob
                        
                        # 使用PIL解析圖片
                        pil_image = Image.open(io.BytesIO(image_data))
                        
                        # 檢查尺寸
                        if pil_image.size[0] < self.min_size[0] or pil_image.size[1] < self.min_size[1]:
                            continue
                        
                        # 生成哈希並檢查重複
                        image_hash = self._calculate_image_hash(image_data)
                        if image_hash in self.extracted_hashes:
                            continue
                        
                        self.extracted_hashes.add(image_hash)
                        
                        # 轉換為標準PNG格式
                        png_data = self._convert_to_png(pil_image)
                        
                        # 生成文件名
                        import time
                        timestamp = int(time.time())
                        filename = f"{document_name}_word_img{img_index:03d}_{timestamp}.png"
                        
                        # 創建提取結果
                        extracted_image = ExtractedImage(
                            image_data=png_data,
                            format="PNG",
                            size=pil_image.size,
                            source_page=1,  # Word文檔沒有頁面概念，使用1
                            source_position=BoundingBox(0, 0, pil_image.size[0], pil_image.size[1]),
                            hash=image_hash,
                            filename=filename,
                            metadata={
                                "source_type": "word_embedded",
                                "original_format": pil_image.format or "unknown",
                                "target_ref": rel.target_ref
                            }
                        )
                        
                        extracted_images.append(extracted_image)
                        img_index += 1
                        
                    except Exception as e:
                        logger.warning(f"Word圖片處理失敗 - 索引{img_index}: {e}")
                        continue
        
        except Exception as e:
            logger.error(f"Word圖片提取失敗: {e}")
            raise
        
        logger.info(f"Word圖片提取完成 - 提取了{len(extracted_images)}張圖片")
        return extracted_images
    
    def extract_from_powerpoint(self, ppt_path: str, document_name: str) -> List[ExtractedImage]:
        """
        從PowerPoint文檔提取圖片
        
        Args:
            ppt_path: PowerPoint文件路徑
            document_name: 文檔名稱
            
        Returns:
            提取的圖片列表
        """
        logger.info(f"開始從PowerPoint文檔提取圖片: {ppt_path}")
        extracted_images = []
        
        try:
            from pptx import Presentation
            
            prs = Presentation(ppt_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                img_index = 0
                
                for shape in slide.shapes:
                    try:
                        # 檢查是否為圖片
                        if hasattr(shape, 'image'):
                            # 獲取圖片數據
                            image_data = shape.image.blob
                            
                            # 使用PIL解析圖片
                            pil_image = Image.open(io.BytesIO(image_data))
                            
                            # 檢查尺寸
                            if pil_image.size[0] < self.min_size[0] or pil_image.size[1] < self.min_size[1]:
                                continue
                            
                            # 生成哈希並檢查重複
                            image_hash = self._calculate_image_hash(image_data)
                            if image_hash in self.extracted_hashes:
                                continue
                            
                            self.extracted_hashes.add(image_hash)
                            
                            # 轉換為標準PNG格式
                            png_data = self._convert_to_png(pil_image)
                            
                            # 獲取形狀位置
                            bbox = BoundingBox(
                                shape.left, shape.top,
                                shape.left + shape.width, 
                                shape.top + shape.height
                            )
                            
                            # 生成文件名
                            import time
                            timestamp = int(time.time())
                            filename = f"{document_name}_slide{slide_num:03d}_img{img_index:03d}_{timestamp}.png"
                            
                            # 創建提取結果
                            extracted_image = ExtractedImage(
                                image_data=png_data,
                                format="PNG",
                                size=pil_image.size,
                                source_page=slide_num,
                                source_position=bbox,
                                hash=image_hash,
                                filename=filename,
                                metadata={
                                    "source_type": "ppt_shape",
                                    "original_format": pil_image.format or "unknown",
                                    "shape_type": str(shape.shape_type),
                                    "slide_title": slide.shapes.title.text if slide.shapes.title else ""
                                }
                            )
                            
                            extracted_images.append(extracted_image)
                            img_index += 1
                            
                    except Exception as e:
                        logger.warning(f"PowerPoint圖片處理失敗 - 幻燈片{slide_num}, 形狀{img_index}: {e}")
                        continue
        
        except Exception as e:
            logger.error(f"PowerPoint圖片提取失敗: {e}")
            raise
        
        logger.info(f"PowerPoint圖片提取完成 - 提取了{len(extracted_images)}張圖片")
        return extracted_images
    
    def _calculate_image_hash(self, image_data: bytes) -> str:
        """計算圖片內容的哈希值用於去重"""
        return hashlib.md5(image_data).hexdigest()
    
    def _convert_to_png(self, pil_image: Image.Image) -> bytes:
        """將PIL圖片轉換為PNG格式的字節數據"""
        if pil_image.mode in ('RGBA', 'LA', 'P'):
            # 保持透明度
            output = io.BytesIO()
            pil_image.save(output, format='PNG')
        else:
            # 轉換為RGB
            if pil_image.mode != 'RGB':
                pil_image = pil_image.convert('RGB')
            output = io.BytesIO()
            pil_image.save(output, format='PNG')
        
        return output.getvalue()
    
    def clear_duplicate_cache(self):
        """清空去重緩存"""
        self.extracted_hashes.clear()
        logger.debug("圖片去重緩存已清空")
    
    def get_extraction_stats(self) -> Dict[str, Any]:
        """獲取提取統計信息"""
        return {
            "total_extracted": len(self.extracted_hashes),
            "supported_formats": self.supported_formats,
            "min_size": self.min_size,
            "max_size": self.max_size
        }
