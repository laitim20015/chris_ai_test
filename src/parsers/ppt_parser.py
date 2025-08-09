"""
PowerPoint演示文稿解析器
PowerPoint Presentation Parser

使用python-pptx解析PowerPoint演示文稿，提取幻燈片內容、圖片和圖表。
嚴格按照項目規格文件的技術選型。

技術特點：
- 原生支持，跨平台兼容
- 幻燈片結構保留
- 演講者備註處理
- 動畫和過渡效果忽略
- 支持.pptx和.ppt格式
"""

from typing import List, Dict, Optional, Union, Any, Tuple
from pathlib import Path
import io
import hashlib
from dataclasses import dataclass

from src.parsers.base import (
    BaseParser, ParsedContent, TextBlock, ImageContent, TableContent,
    DocumentMetadata, ContentType, ImageFormat, ParserError,
    create_bounding_box, generate_content_id
)
from src.config.logging_config import get_logger, log_performance

logger = get_logger("ppt_parser")

class PowerPointParser(BaseParser):
    """PowerPoint演示文稿解析器"""
    
    def __init__(self):
        """初始化PowerPoint解析器"""
        super().__init__("PowerPointParser")
        
        self.supported_formats = ['.pptx', '.ppt']
        
        # 檢查依賴
        try:
            from pptx import Presentation
            self.Presentation = Presentation
            self.available = True
            logger.debug("python-pptx可用")
        except ImportError:
            self.Presentation = None
            self.available = False
            logger.warning("python-pptx不可用，請安裝: pip install python-pptx")
    
    def can_parse(self, file_path: Union[str, Path]) -> bool:
        """
        檢查是否可以解析PowerPoint文件
        
        Args:
            file_path: 文件路徑
            
        Returns:
            bool: 是否可以解析
        """
        if not self.available:
            return False
        
        path = Path(file_path)
        return path.suffix.lower() in self.supported_formats and path.exists()
    
    @log_performance("ppt_parse")
    def parse(self, file_path: Union[str, Path], **kwargs) -> ParsedContent:
        """
        解析PowerPoint演示文稿
        
        Args:
            file_path: PowerPoint文件路徑
            **kwargs: 額外參數
            
        Returns:
            ParsedContent: 解析結果
        """
        if not self.available:
            raise ParserError("python-pptx不可用", parser_name=self.name)
        
        self.validate_file(file_path)
        
        try:
            # 打開PowerPoint演示文稿
            prs = self.Presentation(str(file_path))
            
            text_blocks = []
            images = []
            tables = []
            
            # 逐幻燈片處理
            for slide_num, slide in enumerate(prs.slides, 1):
                # 提取幻燈片文本
                slide_text_blocks = self._extract_slide_text(slide, slide_num)
                text_blocks.extend(slide_text_blocks)
                
                # 提取幻燈片圖片
                slide_images = self._extract_slide_images(slide, slide_num)
                images.extend(slide_images)
                
                # 提取幻燈片表格
                slide_tables = self._extract_slide_tables(slide, slide_num)
                tables.extend(slide_tables)
            
            # 提取元數據
            metadata = self._extract_metadata(prs, Path(file_path))
            metadata.page_count = len(prs.slides)  # 設置幻燈片數量
            
            result = ParsedContent(
                text_blocks=text_blocks,
                images=images,
                tables=tables,
                metadata=metadata,
                success=True
            )
            
            logger.info(f"PowerPoint解析完成: {len(text_blocks)}文本塊, {len(images)}圖片, {len(tables)}表格")
            return result
            
        except Exception as e:
            logger.error(f"PowerPoint解析失敗: {e}")
            return ParsedContent(
                success=False,
                error_message=f"PowerPoint解析失敗: {str(e)}"
            )
    
    def _extract_slide_text(self, slide, slide_num: int) -> List[TextBlock]:
        """提取幻燈片文本"""
        text_blocks = []
        
        try:
            text_index = 0
            
            # 遍歷幻燈片中的所有形狀
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                
                text = shape.text.strip()
                if not text:
                    continue
                
                # 檢測內容類型
                content_type = ContentType.PARAGRAPH
                
                # 檢查是否為標題（通常是第一個文本框或佔位符）
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    if shape.placeholder_format.type == 1:  # 標題佔位符
                        content_type = ContentType.HEADER
                        heading_level = 1
                    elif shape.placeholder_format.type == 2:  # 副標題
                        content_type = ContentType.HEADER
                        heading_level = 2
                    else:
                        heading_level = 0
                else:
                    heading_level = 0
                
                # 獲取位置信息
                bbox = None
                if hasattr(shape, 'left') and hasattr(shape, 'top'):
                    try:
                        bbox = create_bounding_box(
                            shape.left.inches * 72,  # 轉換為點
                            shape.top.inches * 72,
                            shape.width.inches * 72,
                            shape.height.inches * 72
                        )
                    except:
                        bbox = None
                
                text_block = TextBlock(
                    id=generate_content_id("text", slide_num, text_index),
                    content=text,
                    content_type=content_type,
                    bbox=bbox,
                    page_number=slide_num,
                    paragraph_index=text_index,
                    heading_level=heading_level
                )
                
                text_blocks.append(text_block)
                text_index += 1
            
            # 提取演講者備註
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    notes_block = TextBlock(
                        id=generate_content_id("notes", slide_num, 0),
                        content=notes_text,
                        content_type=ContentType.PARAGRAPH,
                        page_number=slide_num,
                        paragraph_index=text_index
                    )
                    text_blocks.append(notes_block)
        
        except Exception as e:
            logger.warning(f"幻燈片 {slide_num} 文本提取失敗: {e}")
        
        return text_blocks
    
    def _extract_slide_images(self, slide, slide_num: int) -> List[ImageContent]:
        """提取幻燈片圖片"""
        images = []
        
        try:
            image_index = 0
            
            for shape in slide.shapes:
                if shape.shape_type == 13:  # 圖片類型
                    try:
                        # 獲取圖片數據
                        image = shape.image
                        image_data = image.blob
                        
                        # 檢測圖片格式
                        if image_data.startswith(b'\xff\xd8'):
                            format_type = ImageFormat.JPEG
                            extension = "jpg"
                        elif image_data.startswith(b'\x89PNG'):
                            format_type = ImageFormat.PNG
                            extension = "png"
                        elif image_data.startswith(b'GIF'):
                            format_type = ImageFormat.GIF
                            extension = "gif"
                        else:
                            format_type = ImageFormat.UNKNOWN
                            extension = "bin"
                        
                        # 獲取位置信息
                        bbox = None
                        try:
                            bbox = create_bounding_box(
                                shape.left.inches * 72,
                                shape.top.inches * 72,
                                shape.width.inches * 72,
                                shape.height.inches * 72
                            )
                        except:
                            bbox = None
                        
                        image_content = ImageContent(
                            id=generate_content_id("image", slide_num, image_index),
                            filename=f"slide_{slide_num:03d}_image_{image_index:03d}.{extension}",
                            format=format_type,
                            data=image_data,
                            bbox=bbox,
                            page_number=slide_num,
                            image_index=image_index
                        )
                        
                        images.append(image_content)
                        image_index += 1
                        
                    except Exception as e:
                        logger.warning(f"幻燈片 {slide_num} 圖片 {image_index} 提取失敗: {e}")
                        continue
        
        except Exception as e:
            logger.warning(f"幻燈片 {slide_num} 圖片提取失敗: {e}")
        
        return images
    
    def _extract_slide_tables(self, slide, slide_num: int) -> List[TableContent]:
        """提取幻燈片表格"""
        tables = []
        
        try:
            table_index = 0
            
            for shape in slide.shapes:
                if shape.shape_type == 19:  # 表格類型
                    try:
                        table = shape.table
                        
                        rows = []
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            rows.append(row_data)
                        
                        if not rows:
                            continue
                        
                        # 檢測表頭
                        headers = []
                        table_rows = rows
                        
                        if len(rows) > 1:
                            # 假設第一行是表頭
                            potential_header = rows[0]
                            if all(cell and cell.strip() for cell in potential_header):
                                headers = potential_header
                                table_rows = rows[1:]
                        
                        # 獲取位置信息
                        bbox = None
                        try:
                            bbox = create_bounding_box(
                                shape.left.inches * 72,
                                shape.top.inches * 72,
                                shape.width.inches * 72,
                                shape.height.inches * 72
                            )
                        except:
                            bbox = None
                        
                        table_content = TableContent(
                            id=generate_content_id("table", slide_num, table_index),
                            rows=table_rows,
                            headers=headers,
                            bbox=bbox,
                            page_number=slide_num,
                            table_index=table_index
                        )
                        
                        tables.append(table_content)
                        table_index += 1
                        
                    except Exception as e:
                        logger.warning(f"幻燈片 {slide_num} 表格 {table_index} 提取失敗: {e}")
                        continue
        
        except Exception as e:
            logger.warning(f"幻燈片 {slide_num} 表格提取失敗: {e}")
        
        return tables
    
    def _extract_metadata(self, prs, file_path: Path) -> DocumentMetadata:
        """提取演示文稿元數據"""
        try:
            core_props = prs.core_properties
            
            metadata = DocumentMetadata(
                filename=file_path.name,
                file_path=str(file_path.absolute()),
                file_size=file_path.stat().st_size,
                file_format=file_path.suffix.lower(),
                title=core_props.title or "",
                author=core_props.author or "",
                subject=core_props.subject or "",
                creator=core_props.author or "",
                created=core_props.created,
                modified=core_props.modified,
                parser_name=self.name,
                parser_version="1.0.0"
            )
            
            # 處理關鍵詞
            if core_props.keywords:
                metadata.keywords = [kw.strip() for kw in core_props.keywords.split(",")]
            
            return metadata
            
        except Exception as e:
            logger.warning(f"元數據提取失敗: {e}")
            return DocumentMetadata(
                filename=file_path.name,
                file_path=str(file_path.absolute()),
                file_size=file_path.stat().st_size,
                file_format=file_path.suffix.lower(),
                parser_name=self.name,
                parser_version="1.0.0"
            )

def parse_powerpoint_slides(file_path: Union[str, Path]) -> ParsedContent:
    """
    便捷的PowerPoint演示文稿解析函數
    
    Args:
        file_path: PowerPoint文件路徑
        
    Returns:
        ParsedContent: 解析結果
    """
    parser = PowerPointParser()
    return parser.parse_with_metadata(file_path)

def extract_slide_content(file_path: Union[str, Path], slide_number: int) -> List[TextBlock]:
    """
    提取指定幻燈片的內容
    
    Args:
        file_path: PowerPoint文件路徑
        slide_number: 幻燈片編號（從1開始）
        
    Returns:
        List[TextBlock]: 文本塊列表
    """
    parser = PowerPointParser()
    result = parser.parse(file_path)
    
    if result.success:
        return [block for block in result.text_blocks if block.page_number == slide_number]
    return []

def extract_ppt_images(file_path: Union[str, Path]) -> List[ImageContent]:
    """
    提取PowerPoint演示文稿中的所有圖片
    
    Args:
        file_path: PowerPoint文件路徑
        
    Returns:
        List[ImageContent]: 圖片列表
    """
    parser = PowerPointParser()
    result = parser.parse(file_path)
    return result.images if result.success else []

